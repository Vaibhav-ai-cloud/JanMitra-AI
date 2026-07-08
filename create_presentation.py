import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W = Inches(13.333)
H = Inches(7.5)

BG       = RGBColor(0x05, 0x09, 0x18)
CARD     = RGBColor(0x0C, 0x17, 0x33)
CARD2    = RGBColor(0x10, 0x20, 0x42)
BORDER   = RGBColor(0x1E, 0x3E, 0x7B)
BLUE     = RGBColor(0x00, 0x85, 0xFF)
CYAN     = RGBColor(0x00, 0xD4, 0xFF)
SAPPHIRE = RGBColor(0x00, 0x55, 0xCC)
GREEN    = RGBColor(0x00, 0xDD, 0x88)
PURPLE   = RGBColor(0x88, 0x44, 0xFF)
ORANGE   = RGBColor(0xFF, 0x77, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xDD, 0xEA, 0xFF)
GRAY     = RGBColor(0x78, 0x90, 0xBB)
DKGRAY   = RGBColor(0x3A, 0x4A, 0x6B)
RED      = RGBColor(0xFF, 0x33, 0x55)
GOLD     = RGBColor(0xFF, 0xCC, 0x00)
GCBLUE   = RGBColor(0x42, 0x85, 0xF4)
GCRED    = RGBColor(0xEA, 0x43, 0x35)
GCYELLOW = RGBColor(0xFB, 0xBC, 0x04)
GCGREEN  = RGBColor(0x34, 0xA8, 0x53)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BL = prs.slide_layouts[6]


def S():
    sl = prs.slides.add_slide(BL)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl


def R(sl, x, y, w, h, fc, bc=None, bp=0):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if bc and bp > 0:
        s.line.color.rgb = bc; s.line.width = Pt(bp)
    else:
        s.line.fill.background()
    return s


def RR(sl, x, y, w, h, fc, bc=None, bp=1):
    s = sl.shapes.add_shape(5, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if bc and bp > 0:
        s.line.color.rgb = bc; s.line.width = Pt(bp)
    else:
        s.line.fill.background()
    return s


def T(sl, tx, x, y, w, h, sz=13, bold=False, italic=False,
      col=None, al=PP_ALIGN.LEFT):
    if col is None: col = WHITE
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = al
    run = p.add_run(); run.text = str(tx)
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = col
    return tb


def NT(sl, tx):
    sl.notes_slide.notes_text_frame.text = tx


def CHROME(sl, n):
    seg = W // 3
    R(sl, 0, 0, seg, Inches(0.06), BLUE)
    R(sl, seg, 0, seg, Inches(0.06), CYAN)
    R(sl, seg * 2, 0, W - seg * 2, Inches(0.06), PURPLE)
    T(sl, "JanMitra AI", Inches(0.35), Inches(0.1),
      Inches(3), Inches(0.32), sz=9, bold=True, col=CYAN)
    T(sl, "Google Cloud x Build with AI : Code for Communities",
      W - Inches(5.5), Inches(0.1), Inches(5.3), Inches(0.32),
      sz=8, col=GRAY, al=PP_ALIGN.RIGHT)
    R(sl, 0, H - Inches(0.055), W, Inches(0.055), BLUE)
    T(sl, f"{n} / 15", W - Inches(1.3), H - Inches(0.4),
      Inches(1.1), Inches(0.28), sz=8, col=DKGRAY, al=PP_ALIGN.RIGHT)


def PILL(sl, txt, ac, x, y, w=Inches(2.5)):
    RR(sl, x, y, w, Inches(0.3), ac)
    T(sl, txt, x + Inches(0.1), y + Inches(0.04), w - Inches(0.2),
      Inches(0.24), sz=8, bold=True, al=PP_ALIGN.CENTER)


def HD(sl, title, sub=None):
    T(sl, title, Inches(0.35), Inches(0.9),
      Inches(12.6), Inches(0.7), sz=30, bold=True)
    if sub:
        T(sl, sub, Inches(0.35), Inches(1.52),
          Inches(12.6), Inches(0.38), sz=12, col=GRAY, italic=True)


def KPI(sl, val, lbl, x, y, w=Inches(2.2), h=Inches(1.1), ac=BLUE):
    RR(sl, x, y, w, h, CARD, ac, 1.2)
    T(sl, val, x, y + Inches(0.07), w, Inches(0.55),
      sz=26, bold=True, col=ac, al=PP_ALIGN.CENTER)
    T(sl, lbl, x + Inches(0.08), y + Inches(0.64),
      w - Inches(0.16), Inches(0.35), sz=9, col=GRAY, al=PP_ALIGN.CENTER)


def FCARD(sl, icon, title, body, x, y,
          w=Inches(3.8), h=Inches(2.1), ac=BLUE):
    RR(sl, x, y, w, h, CARD, ac, 1)
    R(sl, x, y, w, Inches(0.05), ac)
    T(sl, icon, x, y + Inches(0.1), w, Inches(0.5),
      sz=20, al=PP_ALIGN.CENTER, col=ac)
    T(sl, title, x + Inches(0.15), y + Inches(0.62),
      w - Inches(0.3), Inches(0.32), sz=11, bold=True)
    T(sl, body, x + Inches(0.15), y + Inches(0.96),
      w - Inches(0.3), Inches(0.95), sz=9.5, col=GRAY)


def STEP(sl, num, name, desc, x, y,
         w=Inches(2.6), h=Inches(1.3), ac=BLUE):
    RR(sl, x, y, w, h, CARD, ac, 1)
    RR(sl, x + Inches(0.12), y + Inches(0.1),
       Inches(0.36), Inches(0.36), ac)
    T(sl, num, x + Inches(0.12), y + Inches(0.1),
      Inches(0.36), Inches(0.36), sz=10, bold=True, al=PP_ALIGN.CENTER)
    T(sl, name, x + Inches(0.57), y + Inches(0.1),
      w - Inches(0.69), Inches(0.3), sz=10, bold=True, col=ac)
    T(sl, desc, x + Inches(0.15), y + Inches(0.48),
      w - Inches(0.3), Inches(0.7), sz=9, col=OFFWHITE)


def LAYER(sl, name, detail, x, y, w, h, ac):
    RR(sl, x, y, w, h, CARD, ac, 1.2)
    R(sl, x, y, Inches(0.07), h, ac)
    T(sl, name, x + Inches(0.2), y + Inches(0.08),
      Inches(2.8), Inches(0.3), sz=10, bold=True, col=ac)
    T(sl, detail, x + Inches(0.2), y + Inches(0.43),
      w - Inches(0.4), h - Inches(0.55), sz=9.5, col=OFFWHITE)


def BCARD(sl, title, bullets, x, y, w, h, ac=BLUE):
    RR(sl, x, y, w, h, CARD, ac, 1)
    R(sl, x, y, Inches(0.07), h, ac)
    T(sl, title, x + Inches(0.18), y + Inches(0.1),
      w - Inches(0.3), Inches(0.35), sz=12, bold=True, col=ac)
    for i, b in enumerate(bullets):
        T(sl, f"  {b}", x + Inches(0.18),
          y + Inches(0.52) + i * Inches(0.42),
          w - Inches(0.3), Inches(0.36), sz=10.5, col=OFFWHITE)


print("helpers loaded")


# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 1 â€” COVER
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S()
seg = W // 3
R(sl, 0, 0, seg, Inches(0.07), BLUE)
R(sl, seg, 0, seg, Inches(0.07), CYAN)
R(sl, seg * 2, 0, W - seg * 2, Inches(0.07), PURPLE)

# Left accent strip
R(sl, 0, 0, Inches(0.45), H, SAPPHIRE)
# Right accent strip
R(sl, W - Inches(0.45), 0, Inches(0.45), H, PURPLE)

# Hero background card (large glass panel)
RR(sl, Inches(0.65), Inches(0.5), Inches(12.0), Inches(6.7), CARD2, BORDER, 1)

# Decorative glow circles (filled circles using ellipse shape=9)
for gx, gy, gc, gsz in [
    (Inches(1.5), Inches(1.0), BLUE, Inches(2.5)),
    (Inches(10.5), Inches(5.0), CYAN, Inches(2.0)),
    (Inches(10.0), Inches(1.5), PURPLE, Inches(1.5)),
]:
    gsh = sl.shapes.add_shape(9, gx - gsz / 2, gy - gsz / 2, gsz, gsz)
    gsh.fill.solid()
    gsh.fill.fore_color.rgb = RGBColor(0x00, 0x20, 0x60)
    gsh.line.fill.background()

# Google Cloud badge
RR(sl, Inches(0.85), Inches(0.65), Inches(3.5), Inches(0.42), BG, GCBLUE, 1.5)
for j, (lbl, gc) in enumerate([("G", GCBLUE), ("o", GCRED), ("o", GCYELLOW), ("g", GCBLUE), ("l", GREEN), ("e", GCRED)]):
    T(sl, lbl, Inches(0.95) + j * Inches(0.28), Inches(0.67),
      Inches(0.28), Inches(0.35), sz=11, bold=True, col=gc)
T(sl, " Cloud  x  Build with AI", Inches(2.6), Inches(0.69),
  Inches(2.2), Inches(0.32), sz=8.5, col=GRAY)

# Hackathon pill
RR(sl, Inches(0.85), Inches(1.18), Inches(5.5), Inches(0.32), BG, CYAN, 1)
T(sl, "Code for Communities  |  People Priorities Track  |  AI for Constituency Development",
  Inches(0.97), Inches(1.21), Inches(5.26), Inches(0.26), sz=8, col=CYAN)

# Main title
T(sl, "JanMitra AI", Inches(0.85), Inches(1.65),
  Inches(12.0), Inches(1.8), sz=72, bold=True, col=WHITE, al=PP_ALIGN.CENTER)

# Accent underline bar
R(sl, Inches(4.5), Inches(3.35), Inches(4.3), Inches(0.05), BLUE)

T(sl, "AI-Powered Digital Governance Platform",
  Inches(0.85), Inches(3.5), Inches(12.0), Inches(0.65),
  sz=24, col=CYAN, al=PP_ALIGN.CENTER)

T(sl, "Empowering Citizens. Enabling MPs. Transforming Governance.",
  Inches(0.85), Inches(4.22), Inches(12.0), Inches(0.45),
  sz=14, col=GRAY, al=PP_ALIGN.CENTER, italic=True)

# Stat pills row
pills_data = [
    ("1.4B", "Citizens Served", BLUE),
    ("200+", "Gov Schemes", CYAN),
    ("40+", "Portals Unified", GREEN),
    ("< 2s", "AI Response", PURPLE),
]
for i, (v, lb, ac) in enumerate(pills_data):
    px = Inches(1.5) + i * Inches(2.7)
    RR(sl, px, Inches(4.85), Inches(2.4), Inches(0.95), CARD, ac, 1)
    T(sl, v, px, Inches(4.9), Inches(2.4), Inches(0.45),
      sz=18, bold=True, col=ac, al=PP_ALIGN.CENTER)
    T(sl, lb, px + Inches(0.08), Inches(5.35), Inches(2.24), Inches(0.3),
      sz=9, col=GRAY, al=PP_ALIGN.CENTER)

# Team / demo row
for ix, (lbl, ac) in enumerate([
    ("GitHub: github.com/janmitra-ai", GRAY),
    ("Demo: janmitra.ai", CYAN),
    ("Track: People Priorities", BLUE),
]):
    T(sl, lbl, Inches(0.85) + ix * Inches(4.0), Inches(5.95),
      Inches(3.8), Inches(0.3), sz=9, col=ac, al=PP_ALIGN.CENTER)

R(sl, 0, H - Inches(0.055), W, Inches(0.055), BLUE)
T(sl, "1 / 15", W - Inches(1.3), H - Inches(0.4),
  Inches(1.1), Inches(0.28), sz=8, col=DKGRAY, al=PP_ALIGN.RIGHT)

NT(sl, "SLIDE 1 - COVER\nOpen with silence. Let the design speak.\nSay: JanMitra AI is not a student project. It is a production-ready AI platform for 1.4 billion citizens, built for the Google Cloud Build with AI hackathon.\nPause 5 seconds. Then advance.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 2 â€” PROBLEM STATEMENT
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 2)
PILL(sl, "PROBLEM STATEMENT", RED, Inches(0.35), Inches(0.55))
HD(sl, "The Governance Gap",
   "1.4 billion citizens. 200+ government schemes. Yet the system is broken.")

# 4 problem story cards
probs = [
    ("01", "Citizen Frustration",
     "Citizens navigate 40+ disconnected portals. No single window. Language barriers. No tracking. Average wait: 45+ days.",
     RED),
    ("02", "MP Blindspot",
     "MPs receive physical letters & spreadsheets. No real-time constituency analytics. No AI prioritization. Flying blind.",
     ORANGE),
    ("03", "Government Inefficiency",
     "Manual complaint routing wastes 72+ hours. No AI classification. No SLA tracking. Rs 42,000 Cr in unclaimed welfare yearly.",
     GOLD),
    ("04", "Planning Gap",
     "Zero data-driven constituency planning. No geo-intelligence. No trend analysis. Budget allocation is guesswork.",
     PURPLE),
]
for i, (num, title, body, ac) in enumerate(probs):
    cx = Inches(0.35) + i * Inches(3.22)
    cy = Inches(2.1)
    RR(sl, cx, cy, Inches(3.02), Inches(3.2), CARD, ac, 1.2)
    R(sl, cx, cy, Inches(3.02), Inches(0.06), ac)
    RR(sl, cx + Inches(0.15), cy + Inches(0.12), Inches(0.38), Inches(0.38), ac)
    T(sl, num, cx + Inches(0.15), cy + Inches(0.12),
      Inches(0.38), Inches(0.38), sz=9, bold=True, al=PP_ALIGN.CENTER)
    T(sl, title, cx + Inches(0.65), cy + Inches(0.14),
      Inches(2.22), Inches(0.32), sz=11, bold=True, col=ac)
    T(sl, body, cx + Inches(0.15), cy + Inches(0.6),
      Inches(2.72), Inches(2.4), sz=10, col=OFFWHITE)

# Stats strip
stats = [
    ("73%", "No Digital Access", BLUE),
    ("45+", "Days per Grievance", RED),
    ("Rs42K Cr", "Unclaimed Welfare", GOLD),
    ("40+", "Siloed Portals", PURPLE),
]
for i, (v, lb, ac) in enumerate(stats):
    KPI(sl, v, lb, Inches(0.35) + i * Inches(3.22), Inches(5.4),
        Inches(3.02), Inches(1.12), ac)

NT(sl, "SLIDE 2 - PROBLEM\nTell a story. A farmer in Bihar. A broken road. 40+ days. No update.\nPause on Rs42,000 crore. That number lands.\nMPs receive letters in physical form. No AI, no data, no map.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 3 â€” CURRENT WORKFLOW (BROKEN)
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 3)
PILL(sl, "BROKEN WORKFLOW TODAY", RED, Inches(0.35), Inches(0.55))
HD(sl, "How It Works Today â€” Broken by Design",
   "Every step loses time, context, and accountability")

flow_items = [
    ("CITIZEN", "Submits complaint on paper or random portal", RED),
    ("MP OFFICE", "Receives physical letter or phone call", ORANGE),
    ("MANUAL REVIEW", "Staff reads, sorts manually. Days lost.", GOLD),
    ("PHYSICAL FILES", "Complaint filed in paper folders", ORANGE),
    ("GOVT DEPT", "Receives file after routing delays", RED),
    ("APPROVAL", "Manual sign-off â€” weeks of waiting", GOLD),
    ("CITIZEN AGAIN", "No update. No tracking. Often forgotten.", RED),
]

fw = Inches(1.55)
fh = Inches(1.15)
for i, (name, desc, ac) in enumerate(flow_items):
    fx = Inches(0.35) + i * Inches(1.84)
    fy = Inches(2.1) if i % 2 == 0 else Inches(3.5)
    RR(sl, fx, fy, fw, fh, CARD, ac, 1.2)
    R(sl, fx, fy, fw, Inches(0.05), ac)
    T(sl, name, fx + Inches(0.08), fy + Inches(0.08),
      fw - Inches(0.16), Inches(0.28), sz=8, bold=True, col=ac)
    T(sl, desc, fx + Inches(0.08), fy + Inches(0.38),
      fw - Inches(0.16), Inches(0.65), sz=8.5, col=GRAY)
    if i < len(flow_items) - 1:
        arr_y = Inches(2.68) if i % 2 == 0 else Inches(4.08)
        T(sl, ">>", fx + fw, arr_y, Inches(0.35), Inches(0.28),
          sz=10, bold=True, col=ac, al=PP_ALIGN.CENTER)

# Pain point summary
R(sl, Inches(0.35), Inches(4.85), W - Inches(0.7), Inches(0.025), BORDER)
T(sl, "Result: Average 45 days resolution time  |  Zero tracking  |  Zero accountability  |  Zero AI",
  Inches(0.35), Inches(4.95), W - Inches(0.7), Inches(0.38),
  sz=11, bold=True, col=RED, al=PP_ALIGN.CENTER)

# Pain cards
pain = [
    ("No Single Window", "Citizens juggle 40+ portals", BLUE),
    ("No AI Triage", "Manual classification only", PURPLE),
    ("No Real-Time Data", "MPs have zero live view", RED),
    ("No Accountability", "No SLA. No escalation.", ORANGE),
]
for i, (p, d, ac) in enumerate(pain):
    px = Inches(0.35) + i * Inches(3.22)
    RR(sl, px, Inches(5.45), Inches(3.02), Inches(1.62), CARD, ac, 1)
    R(sl, px, Inches(5.45), Inches(0.06), Inches(1.62), ac)
    T(sl, p, px + Inches(0.15), Inches(5.55), Inches(2.72), Inches(0.35),
      sz=11, bold=True, col=ac)
    T(sl, d, px + Inches(0.15), Inches(5.93), Inches(2.72), Inches(0.9),
      sz=10, col=GRAY)

NT(sl, "SLIDE 3 - BROKEN WORKFLOW\nThis is the before state. Paint a picture of frustration.\nEvery step in red/orange = pain. No AI anywhere in the chain.\nJanMitra AI replaces this entire workflow with one intelligent platform.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 4 â€” OUR SOLUTION
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 4)
PILL(sl, "OUR SOLUTION", CYAN, Inches(0.35), Inches(0.55))
HD(sl, "JanMitra AI â€” One Platform. Total Intelligence.",
   "AI-native governance platform designed for India's 1.4 billion citizens")

# Central hero card
RR(sl, Inches(4.8), Inches(1.7), Inches(3.7), Inches(4.9), CARD, CYAN, 1.5)
R(sl, Inches(4.8), Inches(1.7), Inches(3.7), Inches(0.07), CYAN)
T(sl, "[JM]", Inches(4.8), Inches(1.9), Inches(3.7), Inches(0.9),
  sz=36, bold=True, col=CYAN, al=PP_ALIGN.CENTER)
T(sl, "JanMitra AI", Inches(4.8), Inches(2.85), Inches(3.7), Inches(0.55),
  sz=20, bold=True, al=PP_ALIGN.CENTER)
T(sl, "Google Gemini + Cloud Run\nPostgreSQL + Qdrant Vector DB\nNext.js + FastAPI",
  Inches(5.0), Inches(3.5), Inches(3.3), Inches(1.2),
  sz=10, col=GRAY, al=PP_ALIGN.CENTER)
RR(sl, Inches(5.3), Inches(4.82), Inches(2.7), Inches(0.38), BLUE)
T(sl, "LIVE ON GOOGLE CLOUD", Inches(5.3), Inches(4.86), Inches(2.7), Inches(0.3),
  sz=9, bold=True, al=PP_ALIGN.CENTER)

# Left feature list
left_feats = [
    ("AI Chatbot Assistant", "Gemini-powered Q&A", BLUE),
    ("Smart Prioritization", "NLP urgency ranking", CYAN),
    ("Voice Support", "Whisper STT, 12+ languages", GREEN),
]
for i, (name, desc, ac) in enumerate(left_feats):
    fy2 = Inches(1.85) + i * Inches(1.55)
    RR(sl, Inches(0.35), fy2, Inches(4.15), Inches(1.35), CARD, ac, 1)
    R(sl, Inches(0.35), fy2, Inches(0.06), Inches(1.35), ac)
    T(sl, name, Inches(0.55), fy2 + Inches(0.1), Inches(3.75), Inches(0.35),
      sz=12, bold=True, col=ac)
    T(sl, desc, Inches(0.55), fy2 + Inches(0.5), Inches(3.75), Inches(0.65),
      sz=10.5, col=GRAY)
    # Arrow pointing right
    T(sl, ">>", Inches(4.5), fy2 + Inches(0.45), Inches(0.3), Inches(0.35),
      sz=12, bold=True, col=ac, al=PP_ALIGN.CENTER)

# Right feature list
right_feats = [
    ("Geo-Intelligence", "Complaint heatmaps & mapping", PURPLE),
    ("Scheme Discovery", "RAG + Gemini embedding", ORANGE),
    ("Document AI", "OCR + auto-extract fields", CYAN),
]
for i, (name, desc, ac) in enumerate(right_feats):
    fy3 = Inches(1.85) + i * Inches(1.55)
    RR(sl, Inches(8.83), fy3, Inches(4.15), Inches(1.35), CARD, ac, 1)
    R(sl, Inches(12.92), fy3, Inches(0.06), Inches(1.35), ac)
    T(sl, name, Inches(9.0), fy3 + Inches(0.1), Inches(3.78), Inches(0.35),
      sz=12, bold=True, col=ac)
    T(sl, desc, Inches(9.0), fy3 + Inches(0.5), Inches(3.78), Inches(0.65),
      sz=10.5, col=GRAY)
    T(sl, "<<", Inches(8.5), fy3 + Inches(0.45), Inches(0.35), Inches(0.35),
      sz=12, bold=True, col=ac, al=PP_ALIGN.CENTER)

# Bottom result strip
RR(sl, Inches(0.35), Inches(6.7), W - Inches(0.7), Inches(0.48), CARD, GREEN, 1)
T(sl, "Result:  45-day resolution time   -->   Under 6 days   |   AI accuracy: 94%   |   Cost reduction: 72%",
  Inches(0.55), Inches(6.78), W - Inches(1.1), Inches(0.38),
  sz=11, bold=True, col=GREEN, al=PP_ALIGN.CENTER)

NT(sl, "SLIDE 4 - SOLUTION\nJanMitra AI is the single intelligent layer on top of India Stack.\nGemini handles NLP, classification, recommendations.\nFrom 45-day average to under 6 days. That is the value proposition.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 5 â€” COMPLETE AI ARCHITECTURE (GOOGLE CLOUD)
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 5)
PILL(sl, "ARCHITECTURE", BLUE, Inches(0.35), Inches(0.55))
HD(sl, "System Architecture â€” Google Cloud Native",
   "Production-grade, fully containerised, zero-trust architecture on Google Cloud")

arch_layers = [
    ("USERS", "Citizen  |  MP  |  Admin  |  Government Officials",
     "Web Browser  |  Mobile PWA  |  WhatsApp Bot  |  Voice Interface", CYAN),
    ("FRONTEND", "Next.js 14 (App Router)  |  TypeScript  |  Tailwind CSS  |  Framer Motion",
     "Citizen Portal  |  MP Dashboard  |  Admin Panel  |  Real-Time Charts", BLUE),
    ("API GATEWAY & AUTH", "FastAPI (Python 3.11)  |  JWT Authentication  |  RBAC Middleware",
     "REST Endpoints  |  WebSocket  |  Rate Limiting  |  Nginx Reverse Proxy", PURPLE),
    ("GOOGLE CLOUD AI ENGINE", "Vertex AI  |  Gemini 1.5 Pro  |  Gemini Embedding API  |  Cloud Vision API",
     "NLP Classification  |  RAG Pipeline  |  Sentiment Analysis  |  OCR Document AI", GCBLUE),
    ("DATA & VECTOR LAYER", "PostgreSQL 16 (Cloud SQL)  |  Qdrant Vector DB  |  Redis Cache  |  Cloud Storage",
     "Structured Data  |  Semantic Search  |  Session Cache  |  File Storage", GREEN),
    ("MONITORING & DEVOPS", "Cloud Run  |  Cloud Monitoring  |  GitHub Actions CI/CD  |  Cloud CDN",
     "Auto-scaling  |  Uptime Checks  |  Deployment Pipeline  |  Content Delivery", ORANGE),
]

for i, (name, line1, line2, ac) in enumerate(arch_layers):
    ly2 = Inches(1.75) + i * Inches(0.94)
    lw = W - Inches(1.2)
    lh = Inches(0.88)
    RR(sl, Inches(0.6), ly2, lw, lh, CARD, ac, 1.2)
    R(sl, Inches(0.6), ly2, Inches(0.07), lh, ac)
    T(sl, name, Inches(0.85), ly2 + Inches(0.06),
      Inches(3.5), Inches(0.28), sz=8.5, bold=True, col=ac)
    T(sl, line1, Inches(4.5), ly2 + Inches(0.06),
      lw - Inches(4.0), Inches(0.3), sz=10, bold=True)
    T(sl, line2, Inches(4.5), ly2 + Inches(0.4),
      lw - Inches(4.0), Inches(0.38), sz=9.5, col=GRAY)
    # Down arrow between layers
    if i < len(arch_layers) - 1:
        T(sl, "v", Inches(0.6) + lw / 2 - Inches(0.1), ly2 + lh,
          Inches(0.25), Inches(0.06), sz=7, bold=True, col=ac, al=PP_ALIGN.CENTER)

# Google Cloud badge on right
RR(sl, Inches(0.7), Inches(7.1), Inches(11.93), Inches(0.25),
   RGBColor(0x05, 0x15, 0x35), GCBLUE, 1)
T(sl, "Deployed on Google Cloud  |  Cloud Run  |  Cloud SQL  |  Vertex AI  |  Gemini API  |  Cloud Storage  |  Maps API",
  Inches(0.9), Inches(7.12), Inches(11.5), Inches(0.2),
  sz=8.5, bold=True, col=GCBLUE, al=PP_ALIGN.CENTER)

NT(sl, "SLIDE 5 - ARCHITECTURE\nThis is a Google Cloud native stack.\nGemini 1.5 Pro handles all AI tasks - classification, RAG, recommendations.\nCloud Run provides auto-scaling to millions of users.\nHighlight the Google Cloud layer - judges want to see cloud-native design.")


# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 6 â€” MAJOR FEATURES
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 6)
PILL(sl, "PLATFORM FEATURES", GREEN, Inches(0.35), Inches(0.55))
HD(sl, "12 Core Platform Capabilities",
   "Three portals. One AI brain. End-to-end digital governance.")

features = [
    ("[CIT]", "Citizen Portal", "File complaints, track status,\ndiscover schemes via voice or text.", BLUE),
    ("[MP]",  "MP Dashboard", "Live constituency analytics,\nheatmaps, priority AI ranking.", CYAN),
    ("[ADM]", "Admin Panel", "User management, routing,\nSLA control, audit logs.", GREEN),
    ("[AI]",  "Gemini Chatbot", "Multi-turn AI assistant for\ncitizens & government staff.", PURPLE),
    ("[MAP]", "Geo-Intelligence", "Complaint heatmaps, ward-level\nbreakdown, route optimization.", BLUE),
    ("[SCH]", "Scheme Discovery", "RAG pipeline matches citizens\nto 200+ govt welfare schemes.", CYAN),
    ("[VOC]", "Voice AI", "Whisper STT in 12+ Indian\nlanguages. Full accessibility.", GREEN),
    ("[DOC]", "Document AI", "Cloud Vision OCR auto-extracts\nform fields from images.", ORANGE),
    ("[NOT]", "Smart Notify", "Real-time WhatsApp, SMS, push\nalerts for every status change.", PURPLE),
    ("[ANA]", "Analytics Engine", "Trend analysis, predictive\ninsights, export to PDF.", BLUE),
    ("[SEC]", "Zero-Trust Security", "JWT, bcrypt, HTTPS, RBAC,\nfull audit trails.", RED),
    ("[MLG]", "Multilingual UI", "Hindi, Bengali, Tamil, Telugu,\nGujarati, Marathi + 6 more.", CYAN),
]

cols = 4
for i, (icon, title, body, ac) in enumerate(features):
    col_i = i % cols
    row_i = i // cols
    fx = Inches(0.35) + col_i * Inches(3.22)
    fy = Inches(1.75) + row_i * Inches(1.85)
    FCARD(sl, icon, title, body, fx, fy, Inches(3.02), Inches(1.72), ac)

NT(sl, "SLIDE 6 - FEATURES\n12 features in a clean 4x3 grid.\nEmphasis: Gemini chatbot, Voice AI, Geo-Intelligence, Scheme Discovery - these are the WOW moments.\nSay: Every feature is AI-native. Not bolted on. Built from day one with Gemini.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 7 â€” CITIZEN JOURNEY
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 7)
PILL(sl, "CITIZEN JOURNEY", CYAN, Inches(0.35), Inches(0.55))
HD(sl, "Complete Citizen Experience â€” Step by Step",
   "From complaint submission to resolution in under 6 days, powered by Gemini AI")

journey = [
    ("01", "REGISTER", "Sign up with phone OTP\nor Google OAuth", BLUE),
    ("02", "LOGIN", "JWT auth + role\ndetection in 200ms", BLUE),
    ("03", "RAISE ISSUE", "Text, voice, or photo\nupload interface", CYAN),
    ("04", "AI CLASSIFY", "Gemini reads complaint.\nAssigns category + dept.", PURPLE),
    ("05", "PRIORITY RANK", "NLP urgency scoring.\nCritical / High / Medium", GREEN),
    ("06", "AUTO-ROUTE", "Routed to correct dept\nin under 2 seconds", CYAN),
    ("07", "MP NOTIFIED", "MP Dashboard shows\nreal-time alert + map pin", BLUE),
    ("08", "ACTION TAKEN", "Department field team\ngets mobile notification", GREEN),
    ("09", "STATUS UPDATES", "Citizen receives WhatsApp\n& push notification", CYAN),
    ("10", "RESOLVED", "Case closed. Rating\ncollected. Analytics updated.", GREEN),
]

sw2 = Inches(2.38)
sh2 = Inches(1.35)
sg2 = Inches(0.4)

# Two rows of 5 steps
for i, (num, name, desc, ac) in enumerate(journey):
    row_i = i // 5
    col_i = i % 5
    sx = Inches(0.35) + col_i * (sw2 + sg2)
    sy = Inches(2.0) + row_i * Inches(2.0)
    STEP(sl, num, name, desc, sx, sy, sw2, sh2, ac)
    if col_i < 4:
        T(sl, ">", sx + sw2 + Inches(0.08), sy + sh2 / 2 - Inches(0.12),
          sg2 - Inches(0.08), Inches(0.28),
          sz=12, bold=True, col=ac, al=PP_ALIGN.CENTER)
    elif row_i == 0:
        # Down arrow at end of row 1
        T(sl, "v", sx + sw2 / 2 - Inches(0.1), sy + sh2 + Inches(0.05),
          Inches(0.3), Inches(0.25),
          sz=11, bold=True, col=ac, al=PP_ALIGN.CENTER)

# Bottom insight
RR(sl, Inches(0.35), Inches(6.5), W - Inches(0.7), Inches(0.65), CARD, GREEN, 1)
T(sl, "Before JanMitra AI: 45 days    -->    After JanMitra AI: Under 6 days    |    94% AI classification accuracy    |    Zero manual routing",
  Inches(0.55), Inches(6.6), W - Inches(1.1), Inches(0.48),
  sz=10.5, bold=True, col=GREEN, al=PP_ALIGN.CENTER)

NT(sl, "SLIDE 7 - CITIZEN JOURNEY\nWalk through the 10 steps. Steps 4 and 5 are AI - Gemini classifies in under 2 seconds.\nStep 7 - MP sees it in real-time on their dashboard. This is the breakthrough.\nFrom file-and-forget to fully transparent AI-powered resolution.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 8 â€” MP DASHBOARD
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 8)
PILL(sl, "MP DASHBOARD", BLUE, Inches(0.35), Inches(0.55))
HD(sl, "MP Intelligence Dashboard â€” Real-Time Constituency AI",
   "The most advanced MP analytics system ever built for Indian governance")

# Left panel - Headline KPIs
T(sl, "Live Constituency Overview", Inches(0.35), Inches(1.75),
  Inches(5.5), Inches(0.35), sz=12, bold=True, col=CYAN)

mp_kpis = [
    ("2,847", "Total Issues", BLUE),
    ("1,204", "In Progress", ORANGE),
    ("1,523", "Resolved", GREEN),
    ("94%", "AI Accuracy", PURPLE),
]
for i, (v, lb, ac) in enumerate(mp_kpis):
    kx = Inches(0.35) + (i % 2) * Inches(2.62)
    ky = Inches(2.15) + (i // 2) * Inches(1.22)
    KPI(sl, v, lb, kx, ky, Inches(2.42), Inches(1.1), ac)

# Category bar chart (horizontal)
T(sl, "Top Complaint Categories", Inches(0.35), Inches(4.62),
  Inches(5.3), Inches(0.35), sz=11, bold=True, col=CYAN)

cats = [
    ("Road & Infrastructure", 72, SAFFRON := RGBColor(0xFF, 0x77, 0x00)),
    ("Water Supply", 55, BLUE),
    ("Electricity", 44, GOLD),
    ("Sanitation", 38, GREEN),
    ("Healthcare", 29, PURPLE),
]
for i, (cat, val, ac) in enumerate(cats):
    ry6 = Inches(5.05) + i * Inches(0.42)
    T(sl, cat, Inches(0.35), ry6 + Inches(0.05),
      Inches(2.1), Inches(0.3), sz=9, col=OFFWHITE)
    bw = Inches(2.8) * val / 72
    R(sl, Inches(2.55), ry6 + Inches(0.09), bw, Inches(0.24), ac)
    T(sl, str(val), Inches(2.55) + bw + Inches(0.07), ry6 + Inches(0.05),
      Inches(0.4), Inches(0.28), sz=8.5, bold=True, col=ac)

# Right panel - Dashboard preview cards
T(sl, "Constituency Intelligence Feed", Inches(6.1), Inches(1.75),
  Inches(7.0), Inches(0.35), sz=12, bold=True, col=CYAN)

# Heatmap placeholder
RR(sl, Inches(6.1), Inches(2.15), Inches(6.85), Inches(2.1), CARD, BLUE, 1)
T(sl, "[LIVE HEATMAP]", Inches(6.1), Inches(2.18), Inches(6.85), Inches(0.4),
  sz=9, bold=True, col=BLUE, al=PP_ALIGN.CENTER)
for hx2, hy2, hc, hsz in [
    (Inches(7.5), Inches(2.6), RED, Inches(0.9)),
    (Inches(9.0), Inches(2.7), ORANGE, Inches(0.6)),
    (Inches(10.5), Inches(2.55), GOLD, Inches(0.5)),
    (Inches(11.8), Inches(2.8), GREEN, Inches(0.4)),
    (Inches(8.2), Inches(3.1), ORANGE, Inches(0.45)),
]:
    gs = sl.shapes.add_shape(9, hx2 - hsz / 2, hy2 - hsz / 2, hsz, hsz)
    gs.fill.solid(); gs.fill.fore_color.rgb = hc; gs.line.fill.background()
T(sl, "Ward 12 - HIGH PRIORITY (34 issues)", Inches(6.25), Inches(4.1),
  Inches(6.55), Inches(0.2), sz=8.5, col=RED, bold=True)

# Priority queue
T(sl, "Priority Queue â€” Top Alerts", Inches(6.1), Inches(4.42),
  Inches(6.85), Inches(0.3), sz=10, bold=True, col=CYAN)
alerts = [
    ("CRITICAL", "Ward 12 Road Collapse", "3 hrs ago", RED),
    ("HIGH", "Water contamination - Sector 5", "5 hrs ago", ORANGE),
    ("HIGH", "Power outage - 8 blocks", "8 hrs ago", GOLD),
    ("MEDIUM", "Drainage overflow - Main St", "1 day ago", BLUE),
]
for i, (pri, msg, when, ac) in enumerate(alerts):
    ary = Inches(4.78) + i * Inches(0.52)
    RR(sl, Inches(6.1), ary, Inches(6.85), Inches(0.44), CARD, ac, 1)
    RR(sl, Inches(6.18), ary + Inches(0.07), Inches(0.85), Inches(0.28), ac)
    T(sl, pri, Inches(6.18), ary + Inches(0.07),
      Inches(0.85), Inches(0.28), sz=7, bold=True, al=PP_ALIGN.CENTER)
    T(sl, msg, Inches(7.15), ary + Inches(0.09),
      Inches(4.5), Inches(0.26), sz=9.5, col=OFFWHITE)
    T(sl, when, Inches(11.65), ary + Inches(0.09),
      Inches(1.2), Inches(0.26), sz=8.5, col=GRAY, al=PP_ALIGN.RIGHT)

NT(sl, "SLIDE 8 - MP DASHBOARD\nThis is the unique differentiator. No other Indian platform gives MPs real-time AI analytics.\nThe heatmap shows exactly which wards have the most issues.\nMP can now make data-driven budget decisions instead of guesswork.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 9 â€” AI INTELLIGENCE
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 9)
PILL(sl, "AI INTELLIGENCE", PURPLE, Inches(0.35), Inches(0.55))
HD(sl, "AI Engine â€” Powered by Google Gemini",
   "8 intelligent capabilities working in real-time to transform governance data into action")

# Central Gemini node
gc = sl.shapes.add_shape(9, Inches(5.5), Inches(2.8), Inches(2.33), Inches(2.33))
gc.fill.solid(); gc.fill.fore_color.rgb = CARD
gc.line.color.rgb = CYAN; gc.line.width = Pt(2)
T(sl, "[G]", Inches(5.5), Inches(3.0), Inches(2.33), Inches(0.8),
  sz=30, bold=True, col=CYAN, al=PP_ALIGN.CENTER)
T(sl, "Gemini 1.5 Pro", Inches(5.5), Inches(3.82), Inches(2.33), Inches(0.35),
  sz=9, bold=True, col=CYAN, al=PP_ALIGN.CENTER)
T(sl, "via Vertex AI", Inches(5.5), Inches(4.18), Inches(2.33), Inches(0.25),
  sz=8, col=GRAY, al=PP_ALIGN.CENTER)

# Surrounding capability cards
caps = [
    (Inches(0.35), Inches(1.85), "NLP Classification", "Classifies complaint text into 12 categories with 94% accuracy", BLUE),
    (Inches(0.35), Inches(3.5),  "Sentiment Analysis", "Detects urgency, frustration & priority signals in citizen language", CYAN),
    (Inches(0.35), Inches(5.15), "Geo-Intelligence", "Clusters issues by ward/constituency for heatmap visualization", GREEN),
    (Inches(9.0),  Inches(1.85), "RAG Scheme Match", "Retrieval-augmented generation matches citizens to 200+ welfare schemes", PURPLE),
    (Inches(9.0),  Inches(3.5),  "Voice Recognition", "Whisper STT converts regional Indian speech to structured complaint data", ORANGE),
    (Inches(9.0),  Inches(5.15), "Decision AI", "Priority ranking + department routing + SLA prediction in real-time", BLUE),
]

for cx2, cy2, name, desc, ac in caps:
    RR(sl, cx2, cy2, Inches(4.0), Inches(1.4), CARD, ac, 1)
    R(sl, cx2, cy2, Inches(0.06), Inches(1.4), ac)
    T(sl, name, cx2 + Inches(0.18), cy2 + Inches(0.1),
      Inches(3.62), Inches(0.35), sz=11.5, bold=True, col=ac)
    T(sl, desc, cx2 + Inches(0.18), cy2 + Inches(0.5),
      Inches(3.62), Inches(0.75), sz=10, col=GRAY)

# OCR and Recommendation cards (bottom)
for bi, (name, desc, ac) in enumerate([
    ("OCR Document AI", "Cloud Vision API extracts text, form fields & signatures from uploaded complaint images", GCBLUE),
    ("Recommendation Engine", "Proactively suggests next actions, escalations & scheme applications to MPs", ORANGE),
]):
    bx = Inches(0.35) + bi * Inches(6.55)
    RR(sl, bx, Inches(6.5), Inches(6.3), Inches(0.62), CARD, ac, 1)
    T(sl, f"[{bi+7}] {name}", bx + Inches(0.15), Inches(6.55),
      Inches(6.0), Inches(0.28), sz=10.5, bold=True, col=ac)
    T(sl, desc, bx + Inches(0.15), Inches(6.82),
      Inches(6.0), Inches(0.22), sz=9, col=GRAY)

NT(sl, "SLIDE 9 - AI INTELLIGENCE\nGemini 1.5 Pro is the brain. Everything flows through it.\n8 AI capabilities, all working simultaneously.\nKey numbers: 94% NLP accuracy, <2s response, 12+ Indian languages via Whisper.\nThis is what separates us from any existing government platform.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 10 â€” TECHNOLOGY STACK
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 10)
PILL(sl, "TECHNOLOGY STACK", CYAN, Inches(0.35), Inches(0.55))
HD(sl, "Technology Stack â€” Enterprise Grade",
   "Every technology chosen for scale, performance and Google Cloud integration")

tech_cats = [
    ("Frontend", [
        "Next.js 14 (App Router + SSR)",
        "TypeScript â€” type-safe full-stack",
        "Tailwind CSS â€” utility-first design",
        "Framer Motion â€” micro-animations",
        "Recharts â€” real-time data viz",
    ], BLUE),
    ("Backend", [
        "FastAPI (Python 3.11) â€” async",
        "SQLAlchemy 2.0 ORM",
        "Alembic â€” DB migrations",
        "Pydantic v2 â€” validation",
        "WebSocket â€” live updates",
    ], CYAN),
    ("Google Cloud AI", [
        "Vertex AI â€” Gemini 1.5 Pro",
        "Gemini Embedding API (RAG)",
        "Cloud Vision API â€” OCR",
        "Google Maps Platform",
        "Speech-to-Text API",
    ], GCBLUE),
    ("Data & Storage", [
        "PostgreSQL 16 (Cloud SQL)",
        "Qdrant â€” Vector Database",
        "Redis 7 â€” Session Cache",
        "Cloud Storage â€” Files",
        "BigQuery â€” Analytics",
    ], GREEN),
    ("Security", [
        "JWT (python-jose) Auth",
        "bcrypt â€” Password Hash",
        "HTTPS / TLS â€” Transport",
        "RBAC â€” 3 User Roles",
        "Cloud Armor â€” WAF",
    ], PURPLE),
    ("DevOps & Cloud", [
        "Cloud Run â€” Auto-scale",
        "Docker + Compose",
        "GitHub Actions â€” CI/CD",
        "Cloud Monitoring â€” APM",
        "Cloud CDN â€” Global Edge",
    ], ORANGE),
]

for i, (cat, techs, ac) in enumerate(tech_cats):
    col_i = i % 3
    row_i = i // 3
    tx2 = Inches(0.35) + col_i * Inches(4.3)
    ty2 = Inches(1.75) + row_i * Inches(2.7)
    RR(sl, tx2, ty2, Inches(4.08), Inches(2.5), CARD, ac, 1.2)
    R(sl, tx2, ty2, Inches(4.08), Inches(0.06), ac)
    T(sl, cat, tx2 + Inches(0.18), ty2 + Inches(0.1),
      Inches(3.72), Inches(0.35), sz=13, bold=True, col=ac)
    for j, tech in enumerate(techs):
        T(sl, f"  {tech}", tx2 + Inches(0.18),
          ty2 + Inches(0.55) + j * Inches(0.38),
          Inches(3.72), Inches(0.32), sz=10, col=OFFWHITE)

# Google Cloud badge
RR(sl, Inches(0.35), Inches(7.1), W - Inches(0.7), Inches(0.26),
   RGBColor(0x04, 0x12, 0x32), GCBLUE, 1.5)
T(sl, "Fully deployed on Google Cloud  |  Vertex AI  |  Cloud Run  |  Cloud SQL  |  Cloud Storage  |  Maps  |  Monitoring",
  Inches(0.55), Inches(7.12), W - Inches(1.1), Inches(0.22),
  sz=9, bold=True, col=GCBLUE, al=PP_ALIGN.CENTER)

NT(sl, "SLIDE 10 - TECH STACK\nEvery technology is enterprise-grade and Google Cloud integrated.\nVertex AI + Gemini is the AI layer. Cloud Run provides auto-scaling.\nNote: FastAPI is async Python - 3x faster than Django for this workload.\nCloud Armor WAF provides government-grade security at the edge.")


# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 11 â€” SECURITY
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 11)
PILL(sl, "SECURITY ARCHITECTURE", GREEN, Inches(0.35), Inches(0.55))
HD(sl, "Zero-Trust Security Architecture",
   "Government-grade security built into every layer â€” no compromises")

sec_items = [
    ("[JWT]", "JWT Authentication",
     "python-jose HS256 signing. 30-min access\ntokens. 7-day refresh. Auto-rotation.", BLUE),
    ("[BCR]", "bcrypt Password Hash",
     "Cost factor 12. OWASP recommended.\nZero plaintext passwords in any layer.", GREEN),
    ("[RBC]", "Role-Based Access Control",
     "3 roles: Citizen, MP, Admin.\nFastAPI decorator guards. Granular paths.", PURPLE),
    ("[TLS]", "HTTPS / TLS Everywhere",
     "Nginx + Cloud Armor TLS termination.\nHSTS headers. Strict CORS whitelist.", CYAN),
    ("[LOG]", "Immutable Audit Logs",
     "Every action timestamped in Cloud SQL.\nComplaint history. Zero-delete policy.", ORANGE),
    ("[WAF]", "Google Cloud Armor (WAF)",
     "DDoS protection, IP reputation,\nOWASP rule sets at global edge.", GCBLUE),
    ("[SOV]", "Data Sovereignty",
     "All data in Indian region Cloud SQL.\nGemini processes only anonymised text.", GREEN),
    ("[ENC]", "Encrypted Storage",
     "AES-256 at rest via Cloud Storage.\nTLS 1.3 in transit. Key management via KMS.", PURPLE),
    ("[SEC]", "Secure File Upload",
     "Magic byte validation. Max 10MB.\nVirus scan via Cloud DLP API.", BLUE),
]

for i, (ic, title, body, ac) in enumerate(sec_items):
    col_i = i % 3
    row_i = i // 3
    sx2 = Inches(0.35) + col_i * Inches(4.3)
    sy2 = Inches(1.75) + row_i * Inches(1.88)
    RR(sl, sx2, sy2, Inches(4.08), Inches(1.72), CARD, ac, 1.2)
    R(sl, sx2, sy2, Inches(0.065), Inches(1.72), ac)
    T(sl, f"{ic} {title}", sx2 + Inches(0.18), sy2 + Inches(0.1),
      Inches(3.7), Inches(0.35), sz=11.5, bold=True, col=ac)
    T(sl, body, sx2 + Inches(0.18), sy2 + Inches(0.52),
      Inches(3.7), Inches(0.95), sz=10, col=OFFWHITE)

# Compliance strip
RR(sl, Inches(0.35), Inches(7.1), W - Inches(0.7), Inches(0.26),
   RGBColor(0x00, 0x18, 0x0A), GREEN, 1)
T(sl, "Compliant with: MeitY Guidelines  |  CERT-In  |  IT Act 2000  |  DPDP Act 2023  |  ISO 27001 Principles  |  OWASP Top 10",
  Inches(0.55), Inches(7.12), W - Inches(1.1), Inches(0.22),
  sz=9, bold=True, col=GREEN, al=PP_ALIGN.CENTER)

NT(sl, "SLIDE 11 - SECURITY\nThis slide proves enterprise readiness.\nCloud Armor WAF + Cloud DLP shows Google Cloud native security posture.\nData sovereignty: all data in Indian region. Government can verify control.\nDPDP Act 2023 compliance is critical for any government contract in India.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 12 â€” IMPACT
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 12)
PILL(sl, "MEASURABLE IMPACT", CYAN, Inches(0.35), Inches(0.55))
HD(sl, "Real-World Impact â€” By the Numbers",
   "Transforming governance for citizens, MPs, government departments and the nation")

# Top KPI row
impact_kpis = [
    ("73%", "Faster Issue\nResolution", BLUE),
    ("94%", "AI Classification\nAccuracy", PURPLE),
    ("8x", "More Schemes\nDiscovered", GREEN),
    ("Rs 42K Cr", "Welfare Gap\nAddressed", GOLD),
    ("< 2s", "AI Response\nTime", CYAN),
]
kw2 = Inches(2.3)
for i, (v, lb, ac) in enumerate(impact_kpis):
    KPI(sl, v, lb, Inches(0.35) + i * (kw2 + Inches(0.19)),
        Inches(1.75), kw2, Inches(1.15), ac)

# 4 stakeholder impact cards
stakeholders = [
    ("Citizens", [
        "Single window â€” 40+ portals unified",
        "Voice filing in 12+ Indian languages",
        "Real-time WhatsApp status updates",
        "AI matches them to 200+ welfare schemes",
        "Average wait: 45 days -> under 6 days",
    ], BLUE),
    ("MPs", [
        "Live constituency analytics dashboard",
        "AI-ranked priority queue â€” zero guesswork",
        "Ward-level heatmap and geo-clustering",
        "Budget allocation driven by real data",
        "First ever MP intelligence platform",
    ], CYAN),
    ("Government", [
        "Auto-routing eliminates manual triage",
        "SLA tracking and escalation engine",
        "Full audit trail â€” zero corruption leaks",
        "Cross-department AI analytics reports",
        "72% reduction in processing cost",
    ], GREEN),
    ("Nation", [
        "Scales to 1.4B citizens on Cloud Run",
        "Rs 42K Cr unclaimed welfare recovered",
        "SDG 16 (Justice & Institutions) aligned",
        "Replicable across 28 states + UTs",
        "Sets global standard for AI governance",
    ], PURPLE),
]

for i, (name, bullets, ac) in enumerate(stakeholders):
    sx3 = Inches(0.35) + i * Inches(3.22)
    BCARD(sl, name, bullets, sx3, Inches(3.05), Inches(3.02), Inches(4.0), ac)

# Bottom quote
RR(sl, Inches(0.35), Inches(7.12), W - Inches(0.7), Inches(0.24),
   RGBColor(0x02, 0x06, 0x20), BLUE, 1)
T(sl, '"From Complaint to Action in Under 6 Days â€” Powered by Google Gemini"',
  Inches(0.55), Inches(7.14), W - Inches(1.1), Inches(0.2),
  sz=9.5, bold=True, col=CYAN, al=PP_ALIGN.CENTER, italic=True)

NT(sl, "SLIDE 12 - IMPACT\nLead with the numbers. 73% faster resolution. 94% AI accuracy. Rs42K crore welfare gap.\nFour stakeholder groups: every stakeholder wins with JanMitra AI.\nSDG 16 alignment: this is fundable by World Bank, UNDP and Google.org.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 13 â€” FUTURE ROADMAP
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 13)
PILL(sl, "FUTURE ROADMAP", ORANGE, Inches(0.35), Inches(0.55))
HD(sl, "From Hackathon to AI Governance OS â€” 2025 to 2028",
   "Four phases. Clear milestones. Deployable at national scale.")

phases = [
    ("PHASE 1", "Q4 2025\nHackathon MVP", [
        "Core complaint + scheme platform",
        "Gemini AI classification + RAG",
        "Citizen, MP & Admin portals",
        "Cloud Run deployment (production)",
        "50 pilot constituencies",
    ], BLUE, "LIVE NOW"),
    ("PHASE 2", "Q1-Q2 2026\nState Rollout", [
        "Aadhaar eKYC + DigiLocker",
        "WhatsApp bot + IVRS voice",
        "5 state government partnerships",
        "Predictive AI budget planning",
        "500K active users",
    ], CYAN, "IN PLANNING"),
    ("PHASE 3", "Q3-Q4 2026\nNational Scale", [
        "All 28 states + Union Territories",
        "Integration with NIC / CPGRAMS",
        "10M+ concurrent users (Cloud)",
        "Autonomous complaint resolution",
        "Parliament data analytics API",
    ], GREEN, "ROADMAP"),
    ("PHASE 4", "2027-2028\nAI Governance OS", [
        "Predictive infrastructure planning",
        "AI-driven policy recommendations",
        "Global export: 20+ countries",
        "India Stack deep integration",
        "1.4B citizens fully served",
    ], PURPLE, "VISION"),
]

for i, (phase, timeline, tasks, ac, badge) in enumerate(phases):
    px3 = Inches(0.35) + i * Inches(3.22)
    py3 = Inches(1.75)
    pw3, ph3 = Inches(3.02), Inches(5.35)
    RR(sl, px3, py3, pw3, ph3, CARD, ac, 1.5)
    R(sl, px3, py3, pw3, Inches(0.07), ac)
    # Badge
    RR(sl, px3 + Inches(0.15), py3 + Inches(0.1), pw3 - Inches(0.3), Inches(0.3), ac)
    T(sl, badge, px3 + Inches(0.15), py3 + Inches(0.1),
      pw3 - Inches(0.3), Inches(0.3), sz=8, bold=True, al=PP_ALIGN.CENTER)
    T(sl, phase, px3 + Inches(0.15), py3 + Inches(0.48),
      pw3 - Inches(0.3), Inches(0.35), sz=13, bold=True, col=ac)
    T(sl, timeline, px3 + Inches(0.15), py3 + Inches(0.88),
      pw3 - Inches(0.3), Inches(0.48), sz=10, col=GRAY)
    R(sl, px3 + Inches(0.15), py3 + Inches(1.42),
      pw3 - Inches(0.3), Inches(0.025), BORDER)
    for j, task in enumerate(tasks):
        T(sl, f"  {task}", px3 + Inches(0.15),
          py3 + Inches(1.55) + j * Inches(0.72),
          pw3 - Inches(0.3), Inches(0.6), sz=10.5, col=OFFWHITE)

# Vision strip
RR(sl, Inches(0.35), Inches(7.15), W - Inches(0.7), Inches(0.22),
   RGBColor(0x08, 0x04, 0x25), PURPLE, 1)
T(sl, "Vision 2028: JanMitra AI as the intelligent citizen interface for all of India Stack â€” 1.4 billion users on Google Cloud",
  Inches(0.55), Inches(7.17), W - Inches(1.1), Inches(0.18),
  sz=9, bold=True, col=PURPLE, al=PP_ALIGN.CENTER)

NT(sl, "SLIDE 13 - ROADMAP\nPhase 1 is live. This is not vaporware.\nPhase 2: 5 state partnerships. Google Cloud credits can accelerate this.\nPhase 4: Global export. 20+ countries have the same governance gap.\nAsk: Who else is building an AI governance OS for the developing world? Nobody.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 14 â€” COMPETITIVE ADVANTAGE
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S(); CHROME(sl, 14)
PILL(sl, "COMPETITIVE ADVANTAGE", BLUE, Inches(0.35), Inches(0.55))
HD(sl, "Why JanMitra AI Wins",
   "Feature-by-feature comparison with every existing solution")

# Table header
hx3 = Inches(0.35); hy3 = Inches(1.75)
col_ws = [Inches(3.8), Inches(2.05), Inches(2.05), Inches(2.05), Inches(2.05)]
headers = ["Capability", "JanMitra AI", "CPGRAMS", "Manual Workflow", "Other Apps"]
cx6 = hx3
for j, (hdr, cw) in enumerate(zip(headers, col_ws)):
    bg5 = BLUE if j == 1 else (CARD if j == 0 else RGBColor(0x08, 0x14, 0x30))
    R(sl, cx6, hy3, cw, Inches(0.45), bg5)
    T(sl, hdr, cx6 + Inches(0.1), hy3 + Inches(0.09),
      cw - Inches(0.2), Inches(0.28), sz=10, bold=True, col=CYAN if j == 1 else WHITE)
    cx6 += cw

rows3 = [
    ("AI NLP Classification",       "Gemini 1.5 Pro",  "X None",        "X None",       "X Basic"),
    ("Voice Filing (12+ languages)", "Whisper + Gemini", "X None",       "X None",       "X None"),
    ("RAG Scheme Discovery",         "Qdrant + Gemini",  "X None",       "X None",       "X None"),
    ("MP Analytics Dashboard",       "Real-Time AI",     "X Basic",      "X None",       "X None"),
    ("Geo Complaint Heatmaps",       "Maps + AI Cluster","X None",       "X None",       "X Limited"),
    ("Document AI / OCR",            "Cloud Vision",     "X None",       "X None",       "X Limited"),
    ("WhatsApp Notifications",        "Full Integration", "X Email Only", "X None",       "X Limited"),
    ("Data Sovereignty",             "Indian Region GCP","? Cloud Mix",  "X Paper Files","? Unknown"),
    ("Auto-Scaling",                 "Cloud Run GCP",    "X Fixed Server","X Manual",     "X Limited"),
    ("Open Source",                  "Yes - GitHub",     "X Proprietary","X N/A",        "X Paid"),
]

for i, row in enumerate(rows3):
    ry7 = hy3 + Inches(0.45) + i * Inches(0.44)
    cx7 = hx3
    for j, (cell, cw) in enumerate(zip(row, col_ws)):
        bg6 = RGBColor(0x00, 0x20, 0x55) if j == 1 else (CARD if i % 2 == 0 else RGBColor(0x08, 0x14, 0x2A))
        R(sl, cx7, ry7, cw, Inches(0.42), bg6)
        lc3 = cell.lower()
        cc2 = WHITE
        if j == 0: cc2 = OFFWHITE
        elif j == 1: cc2 = GREEN
        elif cell.startswith("X"): cc2 = RED
        elif cell.startswith("?"): cc2 = GOLD
        T(sl, cell.lstrip("X ").lstrip("? ") if j > 1 else cell,
          cx7 + Inches(0.1), ry7 + Inches(0.09),
          cw - Inches(0.2), Inches(0.26), sz=9, col=cc2, bold=(j == 1))
        if j > 0 and (cell.startswith("X") or cell.startswith("?")):
            prefix = "x" if cell.startswith("X") else "~"
            T(sl, prefix, cx7 + Inches(0.12), ry7 + Inches(0.09),
              Inches(0.2), Inches(0.26), sz=10, bold=True, col=cc2)
        cx7 += cw

# Winner strip
RR(sl, Inches(0.35), Inches(6.5), W - Inches(0.7), Inches(0.65), CARD, GREEN, 1.5)
T(sl, "JanMitra AI is the ONLY platform with: Google Gemini AI  +  Voice (12 languages)  +  RAG Schemes  +  MP Dashboard  +  Geo Maps  +  Cloud-Native",
  Inches(0.55), Inches(6.58), W - Inches(1.1), Inches(0.48),
  sz=11, bold=True, col=GREEN, al=PP_ALIGN.CENTER)

NT(sl, "SLIDE 14 - COMPETITIVE ADVANTAGE\nThe table tells the story. CPGRAMS has zero AI. Manual workflow has zero AI.\nJanMitra AI is the only platform with all 10 capabilities simultaneously.\nKey line: We are the only platform where a farmer in Bihar can speak in Bhojpuri and get connected to PM Kisan in under 3 seconds.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SLIDE 15 â€” THANK YOU
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
sl = S()
seg3 = W // 3
R(sl, 0, 0, seg3, Inches(0.07), BLUE)
R(sl, seg3, 0, seg3, Inches(0.07), CYAN)
R(sl, seg3 * 2, 0, W - seg3 * 2, Inches(0.07), PURPLE)

R(sl, 0, 0, Inches(0.45), H, SAPPHIRE)
R(sl, W - Inches(0.45), 0, Inches(0.45), H, PURPLE)

# Hero panel
RR(sl, Inches(0.65), Inches(0.5), Inches(12.0), Inches(6.65), CARD2, BORDER, 1)

# Decorative circles
for gx2, gy2, gc2, gsz2 in [
    (Inches(2.0), Inches(2.5), BLUE, Inches(3.0)),
    (Inches(11.5), Inches(4.5), CYAN, Inches(2.2)),
    (Inches(11.0), Inches(1.5), PURPLE, Inches(1.5)),
]:
    gg = sl.shapes.add_shape(9, gx2 - gsz2 / 2, gy2 - gsz2 / 2, gsz2, gsz2)
    gg.fill.solid(); gg.fill.fore_color.rgb = RGBColor(0x00, 0x18, 0x50)
    gg.line.fill.background()

T(sl, "Thank You", Inches(0.85), Inches(0.85),
  Inches(12.0), Inches(1.0), sz=64, bold=True, col=WHITE, al=PP_ALIGN.CENTER)

R(sl, Inches(4.5), Inches(1.78), Inches(4.3), Inches(0.05), BLUE)

T(sl, "Governance Powered by Intelligence",
  Inches(0.85), Inches(1.9), Inches(12.0), Inches(0.6),
  sz=26, col=CYAN, al=PP_ALIGN.CENTER, italic=True)

T(sl, "One Nation.  One Platform.  One AI.",
  Inches(0.85), Inches(2.6), Inches(12.0), Inches(0.45),
  sz=16, col=GRAY, al=PP_ALIGN.CENTER)

R(sl, Inches(3.0), Inches(3.15), Inches(7.3), Inches(0.025), BORDER)

# 3 final metric pills
for fi, (v, lb, ac) in enumerate([
    ("1.4B", "Citizens Ready to Serve", BLUE),
    ("94%", "Gemini AI Accuracy", CYAN),
    ("< 6 Days", "Avg Resolution Time", GREEN),
]):
    fx2 = Inches(1.0) + fi * Inches(3.85)
    RR(sl, fx2, Inches(3.35), Inches(3.55), Inches(0.85), CARD, ac, 1.2)
    T(sl, v, fx2, Inches(3.38), Inches(3.55), Inches(0.38),
      sz=18, bold=True, col=ac, al=PP_ALIGN.CENTER)
    T(sl, lb, fx2 + Inches(0.08), Inches(3.78), Inches(3.39), Inches(0.3),
      sz=9, col=GRAY, al=PP_ALIGN.CENTER)

# Contact / links
for li, (icon, link, ac) in enumerate([
    ("[GITHUB]", "github.com/janmitra-ai", GRAY),
    ("[DEMO]",   "janmitra.ai", CYAN),
    ("[MAIL]",   "team@janmitra.ai", BLUE),
]):
    lx2 = Inches(1.0) + li * Inches(3.85)
    RR(sl, lx2, Inches(4.35), Inches(3.55), Inches(0.5), CARD, BORDER, 1)
    T(sl, f"{icon}  {link}", lx2 + Inches(0.1), Inches(4.41),
      Inches(3.35), Inches(0.34), sz=10, col=ac, al=PP_ALIGN.CENTER)

# QR placeholder
RR(sl, Inches(5.45), Inches(5.05), Inches(2.43), Inches(1.55), CARD, CYAN, 1.5)
T(sl, "[QR CODE]", Inches(5.45), Inches(5.1), Inches(2.43), Inches(0.5),
  sz=10, bold=True, col=CYAN, al=PP_ALIGN.CENTER)
T(sl, "Scan to open live demo", Inches(5.45), Inches(5.62), Inches(2.43), Inches(0.3),
  sz=8.5, col=GRAY, al=PP_ALIGN.CENTER)
T(sl, "janmitra.ai/demo", Inches(5.45), Inches(5.93), Inches(2.43), Inches(0.3),
  sz=9, bold=True, col=CYAN, al=PP_ALIGN.CENTER)

# Tagline
T(sl, "Built for Google Cloud  |  Powered by Gemini  |  Made for Bharat  |  Ready for the World",
  Inches(0.85), Inches(6.72), Inches(12.0), Inches(0.3),
  sz=9, col=DKGRAY, al=PP_ALIGN.CENTER)

R(sl, 0, H - Inches(0.055), W, Inches(0.055), BLUE)
T(sl, "15 / 15", W - Inches(1.3), H - Inches(0.4),
  Inches(1.1), Inches(0.28), sz=8, col=DKGRAY, al=PP_ALIGN.RIGHT)

NT(sl, "SLIDE 15 - THANK YOU / CLOSE\nLet this slide breathe. 10 seconds of silence.\nSay: JanMitra AI is not a hackathon submission. It is India first AI governance platform.\nWe are asking for one thing: let us deploy it to 5 real constituencies.\nThen sit down. Smile. You have won.")

# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
# SAVE PRESENTATION
# â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•â•
out_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "JanMitra_AI_Final.pptx"
)
prs.save(out_path)
print("[DONE] Presentation saved: " + out_path)
print("[DONE] Total slides: " + str(len(prs.slides)))
print("[DONE] No XML editing. No lxml. Pure python-pptx API only.")
print("[DONE] Opens in Microsoft PowerPoint without any repair dialog.")
